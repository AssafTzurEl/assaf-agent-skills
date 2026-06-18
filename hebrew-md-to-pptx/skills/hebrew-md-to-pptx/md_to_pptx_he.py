#!/usr/bin/env python3
"""
Hebrew slide-Markdown -> PowerPoint (.pptx) converter, RTL-aware.

See SKILL.md for the full input-format spec. Produces a basic, editable .pptx
skeleton (one slide per block, native PowerPoint sections, speaker notes), to be
styled afterwards. Without a template the deck is 16:9 widescreen.

Direction handling
------------------
* Each paragraph gets a base direction (a:pPr/@rtl + @algn) by majority of strong
  directional characters.
* Within a paragraph, text is split into runs by script. Each run is tagged with
  a language (a:rPr/@lang = he-IL or en-US). The language tag (a) makes the
  spell-checker use the right dictionary instead of flagging all Hebrew as typos,
  and (b) makes PowerPoint resolve each run's bidi direction, so neutral
  punctuation (dash, dot, brackets) around English stays on the Hebrew side
  instead of drifting LTR. This mirrors the per-run direction marking the
  hebrew-md-to-docx skill does with w:rtl.
* A single newline inside a paragraph is rendered as a real line break (Markdown
  would otherwise collapse it to a space), so the author's line structure in
  speaker notes is preserved.
* Bulleted/numbered paragraphs get a hanging indent so the marker is spaced from
  the text even in text frames (e.g. notes) that define no list styles.
* Ambiguous (near-even) paragraphs are reported; pin them with --overrides.
* Sections are a Microsoft 2010 extension with no python-pptx API, so the
  <p14:sectionLst> is written directly into presentation.xml.
"""

import argparse
import json
import os
import re
import sys
import uuid
from xml.sax.saxutils import escape

import mistune
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

# 16:9 widescreen (standard PowerPoint), used when no template is supplied.
WIDE_W = Emu(12192000)   # 13.333"
WIDE_H = Emu(6858000)    # 7.5"

# ---------------------------------------------------------------------------
# Direction detection (ported from hebrew-md-to-docx)
# ---------------------------------------------------------------------------

HEB_RE = re.compile(r"[֐-׿יִ-ﭏ]")
LAT_RE = re.compile(r"[A-Za-zÀ-ɏ]")
AMBIGUOUS_BAND = 0.20


def counts(text):
    return len(HEB_RE.findall(text)), len(LAT_RE.findall(text))


def detect_direction(text, base="rtl"):
    h, l = counts(text)
    if h == 0 and l == 0:
        return base, False
    if l == 0:
        return "rtl", False
    if h == 0:
        return "ltr", False
    total = h + l
    direction = "rtl" if h >= l else "ltr"
    return direction, abs(h - l) / total < AMBIGUOUS_BAND


def _char_dir(ch):
    if HEB_RE.match(ch):
        return "rtl"
    if LAT_RE.match(ch) or ch.isdigit():
        return "ltr"
    return "neutral"


def split_runs(text):
    """Split text into (segment, 'rtl'|'ltr') runs under an RTL base. A neutral
    char is LTR only when between two LTR characters; otherwise it joins the RTL
    (Hebrew) flow. Ported from hebrew-md-to-docx."""
    if not text:
        return []
    atoms = []
    cur, cur_cls = [], None
    for ch in text:
        cls = _char_dir(ch)
        if cur_cls is None or cls == cur_cls:
            cur.append(ch); cur_cls = cls
        else:
            atoms.append(["".join(cur), cur_cls]); cur, cur_cls = [ch], cls
    atoms.append(["".join(cur), cur_cls])

    n = len(atoms)
    for i, atom in enumerate(atoms):
        if atom[1] != "neutral":
            continue
        prev_dir = next((atoms[j][1] for j in range(i - 1, -1, -1)
                         if atoms[j][1] != "neutral"), None)
        next_dir = next((atoms[j][1] for j in range(i + 1, n)
                         if atoms[j][1] != "neutral"), None)
        atom[1] = "ltr" if (prev_dir == "ltr" and next_dir == "ltr") else "rtl"

    merged = []
    for seg, d in atoms:
        if merged and merged[-1][1] == d:
            merged[-1][0] += seg
        else:
            merged.append([seg, d])
    return [(s, d) for s, d in merged]


class Director:
    def __init__(self, overrides):
        self.overrides = overrides or {}
        self.idx = 0
        self.ambiguous = []

    def decide(self, text, base="rtl"):
        i = self.idx
        self.idx += 1
        direction, ambiguous = detect_direction(text, base=base)
        ovr = self.overrides.get(str(i))
        if ovr in ("rtl", "ltr"):
            return ovr, i
        if ambiguous:
            self.ambiguous.append((i, direction, text[:60]))
        return direction, i


# ---------------------------------------------------------------------------
# Slide-format parsing
# ---------------------------------------------------------------------------

SECTION_RE = re.compile(r"^#\s+\[\s*Section\b[^\]]*\]\s*$", re.IGNORECASE)
SLIDE_HDR_RE = re.compile(r"^##\s+\[\s*Slide\b[^\]]*\]\s*$", re.IGNORECASE)
SLIDE_LABEL_RE = re.compile(r"^\s*\[\s*Slide\s+\d+[^\]]*\]\s*$", re.IGNORECASE)
SECTION_NAME_RE = re.compile(r"^#\s+\[\s*(.*?)\s*\]\s*$")
TITLE_RE = re.compile(r"^#\s+(.+?)\s*$")
UNTITLED_RE = re.compile(r"^\(\s*(?:slide\b.*?)?untitled.*\)$", re.IGNORECASE)
SEP_NOTES_RE = re.compile(r"^===\s*$")
NOTES_LABEL_RE = re.compile(r"^\s*\*\*\s*Speaker Notes\s*:?\s*\*\*\s*$", re.IGNORECASE)


def section_name(line):
    m = SECTION_NAME_RE.match(line)
    inside = m.group(1) if m else line.strip()
    return re.sub(r"^Section\s*\d*\s*:?\s*", "", inside, flags=re.IGNORECASE).strip()


def extract_brackets(text):
    """Pull presenter instructions ([...]) out of body text -> (clean, [instr]).
    A '[...]' immediately followed by '(' is a Markdown link and kept. Brackets
    kept literally (so '[^]' stays recognisable). One nesting level tolerated."""
    out, instr = [], []
    i, n = 0, len(text)
    while i < n:
        ch = text[i]
        if ch == "[":
            depth, j = 1, i + 1
            while j < n and depth:
                if text[j] == "[":
                    depth += 1
                elif text[j] == "]":
                    depth -= 1
                j += 1
            if depth == 0:
                after = text[j] if j < n else ""
                if after == "(":
                    out.append(text[i:j]); i = j; continue
                instr.append(text[i:j]); i = j; continue
        out.append(ch); i += 1
    return "".join(out), instr


def parse_document(text):
    raw_blocks = re.split(r"(?m)^---\s*$", text)
    slides = []
    for raw in raw_blocks:
        sect, kept = None, []
        for ln in raw.splitlines():
            if SECTION_RE.match(ln):
                sect = section_name(ln); continue
            if SLIDE_HDR_RE.match(ln) or SLIDE_LABEL_RE.match(ln):
                continue
            kept.append(ln)
        body_part, notes_part = _split_notes(kept)
        if not body_part.strip() and not notes_part.strip():
            continue
        title, body_md = _extract_title(body_part)
        body_md, instr = extract_brackets(body_md)
        slides.append({"section": sect, "title": title,
                       "body_md": body_md.strip("\n"),
                       "notes_md": notes_part.strip("\n"),
                       "instructions": instr})
    return slides


def _split_notes(lines):
    for i, ln in enumerate(lines):
        if SEP_NOTES_RE.match(ln):
            notes = lines[i + 1:]
            while notes and not notes[0].strip():
                notes.pop(0)
            if notes and NOTES_LABEL_RE.match(notes[0]):
                notes.pop(0)
            return "\n".join(lines[:i]), "\n".join(notes)
    return "\n".join(lines), ""


def _extract_title(body):
    title, out, found = "", [], False
    for ln in body.splitlines():
        if not found:
            m = TITLE_RE.match(ln)
            if m and not ln.lstrip().startswith("##"):
                cand = m.group(1).strip()
                title = "" if UNTITLED_RE.match(cand) else cand
                found = True
                continue
        out.append(ln)
    return title, "\n".join(out)


# ---------------------------------------------------------------------------
# Markdown helpers
# ---------------------------------------------------------------------------

_md = mistune.create_markdown(renderer=None,
                              plugins=["table", "strikethrough", "url"])


def md_tokens(text):
    return _md(text or "")


def first_inline(text):
    for tok in md_tokens(text):
        if tok.get("type") in ("paragraph", "block_text", "heading"):
            return tok.get("children") or []
    return []


def inline_text(tokens):
    out = []
    for tok in tokens or []:
        t = tok.get("type")
        if t in ("text", "codespan"):
            out.append(tok.get("raw", ""))
        elif t in ("strong", "emphasis", "link", "strikethrough"):
            out.append(inline_text(tok.get("children")))
        elif t in ("softbreak", "linebreak"):
            out.append(" ")
    return "".join(out)


# ---------------------------------------------------------------------------
# Text-frame rendering
# ---------------------------------------------------------------------------

HE_LANG = "he-IL"
EN_LANG = "en-US"
_A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
_BULLET_PRECEDES = ("a:tabLst", "a:defRPr", "a:extLst")
_BULLET_MARGIN = int(Inches(0.3))   # per indent level
_BULLET_HANG = int(Inches(0.25))    # space between marker and text


def set_para_direction(paragraph, rtl):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")
    pPr.set("algn", "r" if rtl else "l")


def set_run_lang(run, lang):
    run._r.get_or_add_rPr().set("lang", lang)


def _set_bullet(paragraph, kind):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    new = []
    if kind == "none":
        new.append(parse_xml(f'<a:buNone {_A}/>'))
    elif kind == "char":
        new.append(parse_xml(f'<a:buFont {_A} typeface="Arial"/>'))
        new.append(parse_xml(f'<a:buChar {_A} char="•"/>'))
    elif kind == "autonum":
        new.append(parse_xml(f'<a:buAutoNum {_A} type="arabicPeriod"/>'))
    if kind in ("char", "autonum"):
        lvl = int(pPr.get("lvl") or 0)
        pPr.set("marL", str(_BULLET_MARGIN * (lvl + 1)))
        pPr.set("indent", str(-_BULLET_HANG))
    anchor = None
    for child in pPr:
        if child.tag in (qn(t) for t in _BULLET_PRECEDES):
            anchor = child; break
    for el in new:
        if anchor is not None:
            anchor.addprevious(el)
        else:
            pPr.append(el)


def add_text(paragraph, text, bold=False, italic=False, url=None):
    """One run per script segment, each tagged with a language so neutral
    punctuation joins the correct side and Hebrew is spell-checked as Hebrew."""
    for seg, d in split_runs(text):
        r = paragraph.add_run(); r.text = seg
        if bold:
            r.font.bold = True
        if italic:
            r.font.italic = True
        set_run_lang(r, HE_LANG if d == "rtl" else EN_LANG)
        if url:
            r.hyperlink.address = url


def add_inline_runs(paragraph, tokens, code_font, bold=False, italic=False):
    for tok in tokens or []:
        t = tok.get("type")
        if t == "text":
            add_text(paragraph, tok.get("raw", ""), bold, italic)
        elif t == "strong":
            add_inline_runs(paragraph, tok.get("children"), code_font, True, italic)
        elif t == "emphasis":
            add_inline_runs(paragraph, tok.get("children"), code_font, bold, True)
        elif t == "strikethrough":
            add_inline_runs(paragraph, tok.get("children"), code_font, bold, italic)
        elif t == "codespan":
            r = paragraph.add_run(); r.text = tok.get("raw", "")
            r.font.name = code_font
            if bold:
                r.font.bold = True
            if italic:
                r.font.italic = True
            set_run_lang(r, EN_LANG)
        elif t == "link":
            url = tok.get("attrs", {}).get("url", "")
            start = len(paragraph.runs)
            add_inline_runs(paragraph, tok.get("children"), code_font, bold, italic)
            for r in paragraph.runs[start:]:
                if url:
                    r.hyperlink.address = url
        elif t in ("linebreak", "softbreak"):
            # Render single newlines as real breaks so the author's line
            # structure (especially in speaker notes) survives.
            paragraph.add_line_break()


class BodyWriter:
    def __init__(self, tf, director, code_font):
        self.tf = tf
        self.first = True
        self.director = director
        self.code_font = code_font

    def _para(self):
        if self.first:
            self.first = False
            return self.tf.paragraphs[0]
        return self.tf.add_paragraph()

    def write(self, tokens):
        for tok in tokens or []:
            t = tok.get("type")
            if t in ("paragraph", "block_text"):
                self._flow(tok, bullet="none")
            elif t == "heading":
                self._flow(tok, bullet="none", bold=True)
            elif t == "list":
                self._list(tok)
            elif t == "block_code":
                self._code(tok)
            elif t == "block_quote":
                self.write(tok.get("children"))

    def _flow(self, tok, bullet, level=0, bold=False):
        text = inline_text(tok.get("children"))
        direction, _ = self.director.decide(text)
        p = self._para()
        p.level = level
        set_para_direction(p, direction == "rtl")
        _set_bullet(p, bullet)
        add_inline_runs(p, tok.get("children"), self.code_font, bold=bold)
        return p

    def _list(self, tok, level=0):
        ordered = tok.get("attrs", {}).get("ordered", False)
        for item in tok.get("children", []):
            first = True
            for child in item.get("children", []):
                ct = child.get("type")
                if ct in ("block_text", "paragraph"):
                    self._flow(child,
                               bullet=("autonum" if ordered else "char") if first else "none",
                               level=level)
                    first = False
                elif ct == "list":
                    self._list(child, level=level + 1)
                else:
                    self.write([child])

    def _code(self, tok):
        for line in tok.get("raw", "").rstrip("\n").split("\n"):
            p = self._para()
            set_para_direction(p, rtl=False)
            _set_bullet(p, "none")
            r = p.add_run(); r.text = line if line else " "
            r.font.name = self.code_font
            set_run_lang(r, EN_LANG)


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------

def collect_tables(tokens):
    return [tok for tok in tokens or [] if tok.get("type") == "table"]


def add_table_shape(slide, tok, left, top, width, director, code_font):
    head, body_rows = None, []
    for child in tok.get("children", []):
        if child["type"] == "table_head":
            head = child["children"]
        elif child["type"] == "table_body":
            body_rows = child["children"]
    ncols = len(head) if head else (len(body_rows[0]["children"]) if body_rows else 0)
    if ncols == 0:
        return top
    nrows = (1 if head else 0) + len(body_rows)
    height = Inches(0.4) * nrows
    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = gf.table
    sample = inline_text((head or [{}])[0].get("children")) if head else ""
    if detect_direction(sample, base="rtl")[0] == "rtl":
        table._tbl.tblPr.set("rtl", "1")
    r = 0
    if head:
        for c, cellfmt in enumerate(head):
            _fill_cell(table.cell(r, c), cellfmt, director, code_font, bold=True)
        r += 1
    for row in body_rows:
        for c, cellfmt in enumerate(row["children"]):
            if c < ncols:
                _fill_cell(table.cell(r, c), cellfmt, director, code_font, bold=False)
        r += 1
    return top + height + Inches(0.2)


def _fill_cell(cell, cellfmt, director, code_font, bold):
    text = inline_text(cellfmt.get("children"))
    direction, _ = director.decide(text)
    p = cell.text_frame.paragraphs[0]
    set_para_direction(p, direction == "rtl")
    add_inline_runs(p, cellfmt.get("children"), code_font, bold=bold)


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

IMAGE_RE = re.compile(r"^\s*image:\s*(.*)$", re.IGNORECASE)


def parse_image(instruction, base_dir):
    inside = instruction.strip()[1:-1].strip()
    m = IMAGE_RE.match(inside)
    if not m:
        return None
    rest = re.sub(r'\s*[—-]\s*".*"\s*$', "", m.group(1)).strip()
    if not rest:
        return None
    path = rest if os.path.isabs(rest) else os.path.join(base_dir, rest)
    return path if os.path.exists(path) else None


def is_hidden_marker(instruction):
    return instruction.strip()[1:-1].strip().lower() == "hidden slide"


def add_image(slide, path, left, top, max_w, max_h):
    from PIL import Image
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 4, 3
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    cx = left + (max_w - w) // 2
    slide.shapes.add_picture(path, cx, top, width=w, height=h)
    return top + h + Inches(0.2)


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------

def write_notes(slide, notes_md, note_instr, director, code_font):
    if not notes_md and not note_instr:
        return
    tf = slide.notes_slide.notes_text_frame
    BodyWriter(tf, director, code_font).write(md_tokens(notes_md))
    if note_instr:
        p = tf.add_paragraph()
        set_para_direction(p, rtl=False)
        _set_bullet(p, "none")
        r = p.add_run(); r.text = "Presenter instructions:"; r.font.bold = True
        set_run_lang(r, EN_LANG)
        for ins in note_instr:
            ip = tf.add_paragraph()
            direction, _ = director.decide(ins)
            set_para_direction(ip, direction == "rtl")
            _set_bullet(ip, "char")
            add_text(ip, ins)


# ---------------------------------------------------------------------------
# Native sections
# ---------------------------------------------------------------------------

SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_sections(prs, slide_sections):
    if not any(slide_sections):
        return
    sld_lst = prs.slides._sldIdLst
    sld_ids = [child.get("id") for child in sld_lst]
    groups, cur_name = [], None
    for sid, sect in zip(sld_ids, slide_sections):
        if sect is not None and (not groups or sect != cur_name):
            cur_name = sect
            groups.append([sect, [sid]])
        elif not groups:
            groups.append(["Default Section", [sid]])
        else:
            groups[-1][1].append(sid)

    parts = [f'<p:extLst xmlns:p="{P_NS}">',
             f'<p:ext uri="{SECTION_EXT_URI}">',
             f'<p14:sectionLst xmlns:p14="{P14_NS}">']
    for name, ids in groups:
        gid = "{" + str(uuid.uuid4()).upper() + "}"
        parts.append(f'<p14:section name="{escape(name, {chr(34): "&quot;"})}" id="{gid}">')
        parts.append("<p14:sldIdLst>")
        for sid in ids:
            parts.append(f'<p14:sldId id="{sid}"/>')
        parts.append("</p14:sldIdLst></p14:section>")
    parts.append("</p14:sectionLst></p:ext></p:extLst>")

    pres_el = sld_lst.getparent()
    existing = pres_el.find(qn("p:extLst"))
    new_ext = parse_xml("".join(parts))
    if existing is not None:
        for ext in new_ext.findall(qn("p:ext")):
            existing.append(ext)
    else:
        pres_el.append(new_ext)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(slides, out_path, template, overrides, code_font, base_dir):
    prs = Presentation(template) if template else Presentation()
    widescreen = template is None
    if widescreen:
        prs.slide_width = WIDE_W
        prs.slide_height = WIDE_H
    director = Director(overrides)
    layout = prs.slide_layouts[1]
    slide_sections, no_title = [], []
    content_w = prs.slide_width - Inches(1.0)   # 0.5" margins each side

    for n, sl in enumerate(slides, start=1):
        slide = prs.slides.add_slide(layout)
        slide_sections.append(sl["section"])

        title_ph = slide.shapes.title
        if title_ph is not None:
            p = title_ph.text_frame.paragraphs[0]
            if sl["title"]:
                direction, _ = director.decide(sl["title"])
                set_para_direction(p, direction == "rtl")
                add_inline_runs(p, first_inline(sl["title"]), code_font)
                if not p.runs:
                    add_text(p, sl["title"])
            else:
                no_title.append(n)

        images, note_instr = [], []
        for ins in sl["instructions"]:
            if is_hidden_marker(ins):
                continue
            img = parse_image(ins, base_dir)
            (images if img else note_instr).append(img or ins)

        body_ph = _body_placeholder(slide)

        # The default layout's placeholders are sized for 4:3; widen them to the
        # widescreen content area, preserving their inherited vertical position
        # (setting only width/left would reset top/height to 0).
        if widescreen:
            for ph in (title_ph, body_ph):
                if ph is None:
                    continue
                top, height = ph.top, ph.height   # capture inherited position
                ph.left = Inches(0.5)
                ph.top = top
                ph.width = content_w
                ph.height = height

        tokens = md_tokens(sl["body_md"])
        tables = collect_tables(tokens)
        flow = [t for t in tokens if t.get("type") != "table"]
        n_par = 0
        if body_ph is not None and flow:
            BodyWriter(body_ph.text_frame, director, code_font).write(flow)
            n_par = len(body_ph.text_frame.paragraphs)

        if body_ph is not None:
            left, width = body_ph.left, body_ph.width
            cursor = body_ph.top
            if flow:
                est = Inches(0.45) * n_par + Inches(0.2)
                cursor = body_ph.top + min(est, int(body_ph.height * 0.6))
        else:
            left, width, cursor = Inches(0.5), content_w, Inches(1.8)

        for tok in tables:
            cursor = add_table_shape(slide, tok, left, cursor, width, director, code_font)
        for img in images:
            avail_h = max(Inches(1.5), (body_ph.top + body_ph.height - cursor)
                          if body_ph is not None else Inches(4))
            cursor = add_image(slide, img, left, cursor, width, avail_h)

        write_notes(slide, sl["notes_md"], note_instr, director, code_font)

    add_sections(prs, slide_sections)
    prs.save(out_path)
    return director, no_title


def _body_placeholder(slide):
    title = slide.shapes.title
    title_el = title._element if title is not None else None
    for ph in slide.placeholders:
        if title_el is not None and ph._element is title_el:
            continue
        return ph
    return None


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="Hebrew slide-Markdown -> PPTX (RTL-aware)")
    ap.add_argument("input")
    ap.add_argument("-o", "--output")
    ap.add_argument("--template", help="base .pptx to inherit theme/layouts from")
    ap.add_argument("--overrides", help="JSON map of paragraph index -> rtl/ltr")
    ap.add_argument("--code-font", default="Consolas")
    args = ap.parse_args()

    out = args.output or re.sub(r"\.md$", "", args.input) + ".pptx"
    overrides = {}
    if args.overrides:
        with open(args.overrides, encoding="utf-8") as f:
            overrides = json.load(f)

    with open(args.input, encoding="utf-8") as f:
        text = f.read()
    slides = parse_document(text)
    if not slides:
        sys.exit("No slides found. Check the input format (--- between slides).")

    base_dir = os.path.dirname(os.path.abspath(args.input))
    director, no_title = build(slides, out, args.template, overrides,
                               args.code_font, base_dir)
    print(f"Wrote {out}  ({len(slides)} slides)")
    if no_title:
        print(f"\nNote: {len(no_title)} slide(s) had no title line: {no_title}",
              file=sys.stderr)
    if director.ambiguous:
        print(f"\n{len(director.ambiguous)} ambiguous paragraph(s) "
              f"(near-even Hebrew/Latin). Pin with --overrides (index -> rtl/ltr):",
              file=sys.stderr)
        for idx, d, snippet in director.ambiguous:
            print(f"  [{idx}] used={d}: {snippet!r}", file=sys.stderr)


if __name__ == "__main__":
    main()

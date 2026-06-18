#!/usr/bin/env python3
"""Extract a .pptx into the slide-Markdown authoring format.

Output is one block per slide, separated by `---`. Within a slide, `===`
separates the body from the speaker notes. Blank lines between blocks are
emitted so Markdown viewers render each as its own paragraph. A typical slide:

    [Section 1: Overview]          (only before the first slide of a section)

    [Slide 3]                      (position label; dropped on re-import)

    # Slide Title

    Body as markdown: bullets, numbered lists, tables, **bold**, *italic*,
    ~~strike~~, `code`, [links](url).

    [image: assets/slide3_img1.png — "alt text"]

    [animation — 2 click steps: ...]

    ===

    **Speaker Notes:**

    Notes as markdown (`*(none)*` when empty).

    ---

Rules:
  * `---` separates slides; `===` separates slide content from speaker notes.
  * `[Section N: Name]` marks a presentation section (emitted only when the deck
    uses sections); `[Slide N]` labels each slide by its position.
  * The first `# heading` is the slide title.
  * Presenter instructions (animations, demos, `[^]` clicks) live in [brackets]
    in the content area; images are emitted as `[image: path — "alt"]`.

The format is documented in full in SKILL.md and in
docs/slide-markdown-format.md at the repo root.

Usage:
    python extract.py input.pptx [-o output.md] [--images-dir DIR] [--no-images]

Notes / limitations are documented in SKILL.md. The timing/animation parser is
best-effort: PowerPoint stores animations as a time-node tree (<p:timing>), which
python-pptx does not expose, so we parse the raw XML. We report effect *category*
(entrance / exit / emphasis / motion) reliably; we do NOT name the specific effect
(Fade vs. Wipe etc.) because that mapping is version-dependent and easy to get wrong.
"""

import argparse
import os
import sys

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# inline run / paragraph formatting
# ----------------------------------------------------------------------------

def shorten(text, limit=70):
    text = " ".join(text.split())
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "\u2026"


def _emph(text, bold, italic, strike, addr):
    """Wrap a piece of text with markdown emphasis/strike/link, preserving outer spaces."""
    if text == "":
        return ""
    core = text.strip()
    if core == "":
        return text  # whitespace only
    lead = text[: len(text) - len(text.lstrip())]
    trail = text[len(text.rstrip()):]
    md = core
    if bold and italic:
        md = f"***{md}***"
    elif bold:
        md = f"**{md}**"
    elif italic:
        md = f"*{md}*"
    if strike:
        md = f"~~{md}~~"
    if addr:
        md = f"[{md}]({addr})"
    return f"{lead}{md}{trail}"


def para_md(para, br_sep=" "):
    """Render a paragraph to inline markdown.

    Walks the <a:p> children so soft line breaks (<a:br>) and fields (<a:fld>,
    e.g. slide number) are handled, and merges adjacent runs that share the same
    formatting so 'Bjarne' + ' Stroustrup' (both italic) becomes one *...* span
    rather than two.

    br_sep controls how <a:br> soft line breaks are rendered: space (default,
    for slide body) or "\\n\\n" (for speaker notes, so Shift+Enter line breaks
    appear as paragraph breaks in Markdown)."""
    runs = list(para.runs)
    ri = 0
    segs = []  # ['run', text, bold, italic, strike, link] or ['br']
    for child in para._p:
        tag = etree.QName(child).localname
        if tag == "r":
            run = runs[ri] if ri < len(runs) else None
            ri += 1
            text = run.text if run is not None else (child.findtext(qn("a:t")) or "")
            if text == "":
                continue
            bold = bool(run.font.bold) if run is not None else False
            italic = bool(run.font.italic) if run is not None else False
            strike = False
            link = None
            if run is not None:
                rPr = run._r.find(qn("a:rPr"))
                if rPr is not None and rPr.get("strike") in ("sngStrike", "dblStrike"):
                    strike = True
                try:
                    link = run.hyperlink.address
                except Exception:
                    link = None
            segs.append(["run", text, bold, italic, strike, link])
        elif tag == "br":
            segs.append(["br"])
        elif tag == "fld":
            text = child.findtext(qn("a:t")) or ""
            if text:
                segs.append(["run", text, False, False, False, None])

    merged = []
    for seg in segs:
        if (seg[0] == "run" and merged and merged[-1][0] == "run"
                and merged[-1][2:] == seg[2:]):
            merged[-1][1] += seg[1]
        else:
            merged.append(list(seg))

    out = []
    for seg in merged:
        if seg[0] == "br":
            out.append(br_sep)  # soft line break -> space (slide body) or newline (notes)
        else:
            out.append(_emph(seg[1], seg[2], seg[3], seg[4], seg[5]))
    return "".join(out).rstrip()


def title_md(shape):
    """One-line title text (collapses soft breaks / multiple paragraphs)."""
    parts = [para_md(p) for p in shape.text_frame.paragraphs]
    return " ".join(" ".join(p.split()) for p in parts if p.strip())


def bullet_kind(para, body_default):
    """Return 'none' | 'num' | 'bullet' based on explicit pPr bullet formatting,
    falling back to body_default when nothing is specified."""
    pPr = para._p.find(qn("a:pPr"))
    if pPr is not None:
        if pPr.find(qn("a:buNone")) is not None:
            return "none"
        if pPr.find(qn("a:buAutoNum")) is not None:
            return "num"
        if pPr.find(qn("a:buChar")) is not None:
            return "bullet"
    return "bullet" if body_default else "none"


def render_text_frame(tf, body_default, double_space=False):
    """Render a text frame to Markdown.

    double_space=True is used for speaker notes: each paragraph is separated
    by a blank line, and soft line breaks (<a:br> / Shift+Enter) are rendered
    as paragraph separators rather than spaces."""
    br_sep = "\n\n" if double_space else " "
    lines = []
    counters = {}
    for para in tf.paragraphs:
        level = para.level or 0
        md = para_md(para, br_sep=br_sep)
        if md.strip() == "":
            lines.append("")
            continue
        kind = bullet_kind(para, body_default)
        indent = "  " * level
        if kind == "none":
            lines.append(f"{indent}{md}")
            counters.clear()
        elif kind == "num":
            n = counters.get(level, 0) + 1
            counters[level] = n
            for deeper in [k for k in counters if k > level]:
                del counters[deeper]
            lines.append(f"{indent}{n}. {md}")
        else:  # bullet
            counters.pop(level, None)
            lines.append(f"{indent}- {md}")
    while lines and lines[0] == "":
        lines.pop(0)
    while lines and lines[-1] == "":
        lines.pop()
    if double_space:
        # Join non-empty paragraphs with blank lines; preserve existing blank
        # lines by collapsing consecutive empties before joining.
        chunks, buf = [], []
        for line in lines:
            if line == "":
                if buf:
                    chunks.append("\n".join(buf))
                    buf = []
            else:
                buf.append(line)
        if buf:
            chunks.append("\n".join(buf))
        return "\n\n".join(chunks)
    return "\n".join(lines)


# ----------------------------------------------------------------------------
# shape helpers
# ----------------------------------------------------------------------------

def alt_text(shape):
    cNvPr = shape._element.find(".//" + qn("p:cNvPr"))
    if cNvPr is None:
        return ""
    descr = cNvPr.get("descr") or ""
    return descr.strip()


def is_body_placeholder(shape):
    if not shape.is_placeholder:
        return False
    try:
        ph = shape.placeholder_format.type
    except Exception:
        return False
    return ph in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)


def shape_label(shape):
    """Short human label used when describing animation targets."""
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    return '"' + shorten(p.text.strip()) + '"'
    except Exception:
        pass
    return shape.name or "shape"


def register(shape, registry, label=None, paragraphs=None):
    try:
        sid = str(shape.shape_id)
    except Exception:
        return
    info = {"label": label or shape_label(shape), "paragraphs": paragraphs}
    if paragraphs is None:
        try:
            if shape.has_text_frame:
                info["paragraphs"] = [p.text for p in shape.text_frame.paragraphs]
        except Exception:
            pass
    registry[sid] = info


def render_table(table):
    rows = list(table.rows)
    if not rows:
        return ""

    def cell_md(cell):
        parts = []
        for p in cell.text_frame.paragraphs:
            t = para_md(p).strip()
            if t:
                parts.append(t)
        return (" ".join(parts)).replace("|", "\\|").strip() or " "

    header = [cell_md(c) for c in rows[0].cells]
    out = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row in rows[1:]:
        out.append("| " + " | ".join(cell_md(c) for c in row.cells) + " |")
    return "\n".join(out)


def graphic_uri(shape):
    gd = shape._element.find(".//" + qn("a:graphicData"))
    return gd.get("uri") if gd is not None else None


# ----------------------------------------------------------------------------
# per-shape dispatch
# ----------------------------------------------------------------------------

def process_shape(shape, ctx):
    """Register the shape and return a list of markdown content blocks."""
    registry = ctx["registry"]
    blocks = []
    st = shape.shape_type

    # group -> recurse
    if st == MSO_SHAPE_TYPE.GROUP:
        register(shape, registry, label=shape.name or "group")
        for child in shape.shapes:
            blocks.extend(process_shape(child, ctx))
        return blocks

    # picture (incl. placeholder pictures)
    is_picture = st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
    if not is_picture:
        try:
            _ = shape.image  # raises if not a picture
            is_picture = True
        except Exception:
            is_picture = False
    if is_picture:
        fname = save_image(shape, ctx)
        alt = alt_text(shape)
        ref = f"{ctx['images_subdir']}/{fname}" if ctx["extract_images"] else fname
        label = f'image "{fname}"'
        register(shape, registry, label=label, paragraphs=None)
        instr = f"[image: {ref}" + (f' \u2014 "{alt}"' if alt else "") + "]"
        blocks.append(instr)
        return blocks

    # table
    if shape.has_table:
        register(shape, registry, label="table")
        md = render_table(shape.table)
        if md:
            blocks.append(md)
        return blocks

    # chart
    if shape.has_chart:
        chart = shape.chart
        try:
            ctype = str(chart.chart_type).split(".")[-1].split(" ")[0]
        except Exception:
            ctype = "chart"
        title = ""
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                title = chart.chart_title.text_frame.text.strip()
        except Exception:
            pass
        register(shape, registry, label="chart")
        blocks.append(f"[chart: {ctype}" + (f' \u2014 "{title}"' if title else "") + "]")
        return blocks

    # SmartArt / diagram (graphicFrame)
    uri = graphic_uri(shape)
    if uri and "diagram" in uri:
        register(shape, registry, label="SmartArt/diagram")
        # best-effort: pull any <a:t> text inside the frame element
        texts = [t.text for t in shape._element.iter(qn("a:t")) if t.text and t.text.strip()]
        if texts:
            joined = "; ".join(shorten(t, 40) for t in texts)
            blocks.append(f"[SmartArt/diagram \u2014 text: {joined}]")
        else:
            blocks.append("[SmartArt/diagram \u2014 text not extractable]")
        return blocks

    # text-bearing shape (placeholder text box, autoshape with text, etc.)
    if shape.has_text_frame and shape.text_frame.text.strip():
        register(shape, registry)
        md = render_text_frame(shape.text_frame, body_default=is_body_placeholder(shape))
        if md.strip():
            blocks.append(md)
        ctx.setdefault("text_lines", []).extend(
            p.text for p in shape.text_frame.paragraphs if p.text.strip())
        return blocks

    # anything else: register so animations can still name it
    register(shape, registry)
    return blocks


def save_image(shape, ctx):
    img = shape.image
    ext = (img.ext or "png").lstrip(".")
    ctx["img_count"] += 1
    fname = f"slide{ctx['idx']}_img{ctx['img_count']}.{ext}"
    if ctx["extract_images"]:
        os.makedirs(ctx["images_dir"], exist_ok=True)
        with open(os.path.join(ctx["images_dir"], fname), "wb") as fh:
            fh.write(img.blob)
        ctx["assets"].append(os.path.join(ctx["images_dir"], fname))
    return fname


# ----------------------------------------------------------------------------
# animation / timing parser (best-effort)
# ----------------------------------------------------------------------------

CLASS_NAMES = {
    "entr": "entrance",
    "exit": "exit",
    "emph": "emphasis",
    "path": "motion path",
    "verb": "media",
    "mediacall": "media",
}


def describe_target(spTgt, registry):
    spid = spTgt.get("spid")
    info = registry.get(str(spid))
    label = info["label"] if info else f"shape #{spid}"
    pRg = spTgt.find(qn("p:txEl") + "/" + qn("p:pRg"))
    if pRg is not None and info and info.get("paragraphs"):
        try:
            st = int(pRg.get("st", 0))
            en = int(pRg.get("end", st))
        except Exception:
            return label
        chunk = [t for t in info["paragraphs"][st: en + 1] if t and t.strip()]
        if chunk:
            return '"' + shorten(" / ".join(chunk)) + '"'
    return label


def parse_effect(eff, registry):
    pc = eff.get("presetClass")
    cname = CLASS_NAMES.get(pc, pc or "animation")
    spTgt = eff.find(".//" + qn("p:spTgt"))
    target = describe_target(spTgt, registry) if spTgt is not None else "shape"
    return f"{cname} of {target}"


def parse_click_group(par, registry):
    cTn = par.find(qn("p:cTn"))
    if cTn is None:
        return None
    trigger = "click"
    cond = cTn.find(qn("p:stCondLst") + "/" + qn("p:cond"))
    if cond is not None:
        delay = cond.get("delay")
        evt = cond.get("evt")
        if delay == "indefinite" or evt == "onClick":
            trigger = "click"
        elif delay and delay.isdigit():
            trigger = "auto"
    effects = []
    for node in cTn.iter(qn("p:cTn")):
        if node.get("presetClass") is not None:
            effects.append(parse_effect(node, registry))
    return {"trigger": trigger, "effects": [e for e in effects if e]}


def parse_buildlist_only(timing, registry):
    notes = []
    for b in timing.iter(qn("p:bldP")):
        spid = b.get("spid")
        build = b.get("build")
        info = registry.get(str(spid))
        label = info["label"] if info else f"shape #{spid}"
        if build in ("p", "byParagraph"):
            notes.append(f"{label}: reveals one paragraph per click")
        else:
            notes.append(f"{label}: animated")
    if notes:
        return "[animation:\n  " + "\n  ".join(notes) + "]"
    return ""


def parse_animations(slide, registry):
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        return ""
    try:
        mainseq = None
        for seq in timing.iter(qn("p:seq")):
            if seq.get("nodeType") == "mainSeq":
                mainseq = seq
                break
            if mainseq is None:
                mainseq = seq
        steps = []
        if mainseq is not None:
            cTn = mainseq.find(qn("p:cTn"))
            child = cTn.find(qn("p:childTnLst")) if cTn is not None else None
            if child is not None:
                for par in child.findall(qn("p:par")):
                    step = parse_click_group(par, registry)
                    if step and step["effects"]:
                        steps.append(step)
        if not steps:
            return parse_buildlist_only(timing, registry)

        lines = []
        click_no = 0
        for s in steps:
            effs = "; ".join(s["effects"]) or "animation"
            if s["trigger"] == "click":
                click_no += 1
                lines.append(f"  click {click_no}: {effs}")
            else:
                lines.append(f"  (with/after previous): {effs}")
        header = f"[animation \u2014 {click_no} click step{'s' if click_no != 1 else ''}:"
        return header + "\n" + "\n".join(lines) + "]"
    except Exception as exc:  # never let animation parsing break extraction
        return f"[animation present \u2014 could not parse ({exc.__class__.__name__})]"


# ----------------------------------------------------------------------------
# slide + presentation assembly
# ----------------------------------------------------------------------------

def extract_slide(slide, idx, ctx):
    registry = {}
    ctx["registry"] = registry
    ctx["idx"] = idx
    ctx["img_count"] = 0
    ctx["text_lines"] = []

    title_shape = slide.shapes.title
    title_text = title_md(title_shape) if title_shape is not None else ""

    content_blocks = []
    for shape in slide.shapes:
        if title_shape is not None and shape._element is title_shape._element:
            register(shape, registry)  # register title for animation lookups
            continue
        content_blocks.extend(process_shape(shape, ctx))

    anim = parse_animations(slide, registry)
    if anim:
        content_blocks.append(anim)

    hidden = slide._element.get("show") == "0"

    head = f"# {title_text}" if title_text else "# (untitled)"
    body_parts = []
    if hidden:
        body_parts.append("[hidden slide]")
    body_parts.extend(b for b in content_blocks if b.strip())
    body = "\n\n".join(body_parts)
    content = f"[Slide {idx}]\n\n" + head + ("\n\n" + body if body else "")

    notes_md = ""
    if slide.has_notes_slide:
        tf = slide.notes_slide.notes_text_frame
        if tf is not None:
            notes_md = render_text_frame(tf, body_default=False, double_space=True).strip()
    notes_section = "**Speaker Notes:**\n\n" + (notes_md if notes_md else "*(none)*")

    meta = {
        "idx": idx,
        "title_raw": title_text,
        "title": _norm(title_text),
        "lines": {_norm(x) for x in ctx["text_lines"] if _norm(x)},
    }
    return content + "\n\n===\n\n" + notes_section, meta


_MD_MARKS = str.maketrans("", "", "*~`")


def _norm(s):
    return " ".join((s or "").translate(_MD_MARKS).split()).lower()


def _related_body(a, b):
    """Do two slides share enough body text to be the same idea continued?

    Uses the overlap coefficient (shared lines / smaller slide's line count) so
    that an additive build (earlier text fully contained in the later slide)
    scores 1.0, while a lightly-reworded build still passes. Slides with no body
    text (e.g. an image-only overlay sequence) fall back to the title match."""
    la, lb = a["lines"], b["lines"]
    if not la or not lb:
        return True  # nothing to compare -> rely on the shared title
    overlap = len(la & lb) / min(len(la), len(lb))
    return overlap >= 0.5


def detect_build_sequences(metas):
    """Find runs of consecutive slides that share a title AND share body text —
    these are almost always a single idea built up across slides (progressive
    disclosure split into multiple physical slides, or one base with changing
    overlays). Requiring body overlap avoids grouping unrelated slides that
    merely happen to share a title."""
    seqs = []
    i, n = 0, len(metas)
    while i < n:
        j = i
        title = metas[i]["title"]
        if title:
            while (j + 1 < n and metas[j + 1]["title"] == title
                   and _related_body(metas[j], metas[j + 1])):
                j += 1
        if j > i:
            additive = all(metas[k]["lines"] <= metas[k + 1]["lines"]
                           for k in range(i, j))
            grew = any(len(metas[k + 1]["lines"]) > len(metas[k]["lines"])
                       for k in range(i, j))
            kind = ("each slide adds to the previous (progressive build)"
                    if additive and grew
                    else "shared base, content/overlay changes per slide")
            disp = shorten(metas[i]["title_raw"], 55) or "(untitled)"
            seqs.append(
                f'Slides {metas[i]["idx"]}\u2013{metas[j]["idx"]} share the title '
                f'"{disp}" \u2014 {kind}.')
        i = j + 1
    return seqs


def get_slide_sections(prs):
    """Return a dict mapping 1-based slide position -> (section_idx, section_name).

    PowerPoint stores sections in an extension namespace (p14). We parse the
    presentation XML directly because python-pptx doesn't expose sections.
    Returns an empty dict if the deck has no sections.

    Section slide lists use the same numeric `id` attribute as the main
    p:sldIdLst (e.g. id="256"), NOT the relationship r:id."""
    NS14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    root = prs.part._element

    # Build numeric slide id -> 1-based position
    sldIdLst = root.find(qn("p:sldIdLst"))
    numid_to_pos = {}
    if sldIdLst is not None:
        for pos, el in enumerate(sldIdLst, start=1):
            num_id = el.get("id")
            if num_id:
                numid_to_pos[num_id] = pos

    ext_lst = root.find(qn("p:extLst"))
    if ext_lst is None:
        return {}

    pos_to_section = {}
    sec_idx = 0
    found = False
    for ext in ext_lst:
        sec_lst = ext.find(f"{{{NS14}}}sectionLst")
        if sec_lst is None:
            continue
        found = True
        for sec in sec_lst.findall(f"{{{NS14}}}section"):
            sec_idx += 1
            name = sec.get("name", f"Section {sec_idx}")
            sl_lst = sec.find(f"{{{NS14}}}sldIdLst")
            if sl_lst is None:
                continue
            for sld_el in sl_lst:
                num_id = sld_el.get("id")
                if num_id and num_id in numid_to_pos:
                    pos_to_section[numid_to_pos[num_id]] = (sec_idx, name)
    return pos_to_section if found else {}


def extract(path, images_dir, extract_images=True, slides_wanted=None):
    prs = Presentation(path)
    ctx = {
        "images_dir": images_dir,
        "images_subdir": os.path.basename(images_dir.rstrip("/")) or "images",
        "extract_images": extract_images,
        "assets": [],
    }
    slide_sections = get_slide_sections(prs)
    blocks, metas = [], []
    current_section = None
    for i, s in enumerate(prs.slides, start=1):
        if slides_wanted is not None and i not in slides_wanted:
            continue
        block, meta = extract_slide(s, i, ctx)
        # Prepend section header when this slide opens a new section
        if slide_sections and i in slide_sections:
            sec_key = slide_sections[i]
            if sec_key != current_section:
                current_section = sec_key
                sec_idx, sec_name = sec_key
                block = f"[Section {sec_idx}: {sec_name}]\n\n" + block
        blocks.append(block)
        metas.append(meta)
    return "\n\n---\n\n".join(blocks), ctx["assets"], detect_build_sequences(metas)


def parse_slide_spec(spec, total):
    """Parse '3-20', '3,5,7-9', etc. into a set of 1-based indices."""
    wanted = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            a = int(a) if a.strip() else 1
            b = int(b) if b.strip() else total
            wanted.update(range(a, b + 1))
        else:
            wanted.add(int(part))
    return {i for i in wanted if 1 <= i <= total}


def main():
    ap = argparse.ArgumentParser(description="Extract a .pptx into slide-Markdown format.")
    ap.add_argument("input", help="path to .pptx file")
    ap.add_argument("-o", "--output", help="output .md path (default: <input stem>.md)")
    ap.add_argument("--images-dir", help="directory for extracted images "
                    "(default: <output stem>_images next to the output file)")
    ap.add_argument("--no-images", action="store_true",
                    help="do not write image files; still note [image: ...] inline")
    ap.add_argument("--slides", help="subset of slides to extract, e.g. '3-20' or "
                    "'1,4,7-9' (1-based; default: all)")
    args = ap.parse_args()

    if not os.path.isfile(args.input):
        sys.exit(f"error: file not found: {args.input}")

    out_path = args.output or (os.path.splitext(os.path.basename(args.input))[0] + ".md")
    out_dir = os.path.dirname(os.path.abspath(out_path))
    stem = os.path.splitext(os.path.basename(out_path))[0]
    images_dir = args.images_dir or os.path.join(out_dir, f"{stem}_images")

    slides_wanted = None
    if args.slides:
        total = len(Presentation(args.input).slides._sldIdLst)
        slides_wanted = parse_slide_spec(args.slides, total)
        if not slides_wanted:
            sys.exit(f"error: --slides '{args.slides}' selects no slides (deck has {total})")

    text, assets, sequences = extract(args.input, images_dir,
                                      extract_images=not args.no_images,
                                      slides_wanted=slides_wanted)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(text + "\n")

    print(f"wrote {out_path}", file=sys.stderr)
    if assets:
        print(f"extracted {len(assets)} image(s) -> {images_dir}", file=sys.stderr)
    if sequences:
        print("\nBUILD SEQUENCES DETECTED (extracted as separate slides; "
              "offer to compact each into one progressive-disclosure slide):",
              file=sys.stderr)
        for s in sequences:
            print("  - " + s, file=sys.stderr)


if __name__ == "__main__":
    main()
